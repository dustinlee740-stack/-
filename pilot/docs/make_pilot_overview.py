# -*- coding: utf-8 -*-
"""
코나카드 영향 분석 도구 — 워크숍용 PPT 생성기.

입력 없음 (본문은 모두 사용자 작성 원문이 이 파일에 내장됨)
출력 D:\\da\\pilot\\docs\\pilot_overview.pptx (27 slides, 16:9)

xlsx_to_md.py / impact.py 와 동일하게 절대경로 하드코딩 + 결정론적 생성기 패턴.
재실행 시 동일 결과.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

OUT = r"D:\da\pilot\docs\pilot_overview.pptx"

# ── 색상 (미니멀: 네이비/그레이/액센트 1색) ─────────────────────────────
NAVY    = RGBColor(0x1F, 0x2A, 0x44)
GRAY    = RGBColor(0x55, 0x5B, 0x66)
GRAY_L  = RGBColor(0x9A, 0xA0, 0xA6)
ACCENT  = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x1E, 0x7E, 0x34)
LIGHT   = RGBColor(0xF4, 0xF5, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0xD0, 0xD3, 0xD8)

FONT_KO = "Pretendard"  # 모던·고가독 한글 폰트 (사용자 시스템에 9 weights 설치 확인)
DEFAULT_LS = 1.25       # 단락 줄 간격 (가독성)

# 슬라이드 크기 (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)

# ── 저수준 헬퍼 ───────────────────────────────────────────────────────────
def _set_run(run, text, *, size=14, bold=False, color=GRAY, font=FONT_KO):
    # 안전장치 — OOXML 텍스트 노드에 raw newline은 PowerPoint가 거부.
    # 호출자가 _split_lines로 단락 분할 못 한 경우엔 공백으로 정리.
    if isinstance(text, str) and ("\n" in text or "\r" in text):
        text = text.replace("\r\n", " ").replace("\n", " ").replace("\r", " ")
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한국어 폰트 명시 적용 (asian font, complex script)
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        ea = rPr.find(qn(f"a:{tag}"))
        if ea is None:
            ea = etree.SubElement(rPr, qn(f"a:{tag}"))
        ea.set("typeface", font)


def _apply_paragraph_style(p, *, align=PP_ALIGN.LEFT, line_spacing=DEFAULT_LS,
                            space_after=2):
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)


def _split_lines(text):
    """문자열에 \\n이 있으면 paragraph 분할용 리스트로 변환.
 PowerPoint는 OOXML 텍스트 노드 안의 raw newline을 안전하게 처리하지 못한다 —
 여러 줄이 필요하면 별도 paragraph로 나눠야 안 깨진다."""
    if isinstance(text, str) and "\n" in text:
        return text.split("\n")
    return text


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False,
              color=GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KO,
              line_spacing=DEFAULT_LS):
    text = _split_lines(text)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    _apply_paragraph_style(p, align=align, line_spacing=line_spacing)
    if isinstance(text, list):
        for i, line in enumerate(text):
            if i > 0:
                p = tf.add_paragraph()
                _apply_paragraph_style(p, align=align, line_spacing=line_spacing)
            run = p.add_run()
            _set_run(run, line, size=size, bold=bold, color=color, font=font)
    else:
        run = p.add_run()
        _set_run(run, text, size=size, bold=bold, color=color, font=font)
    return box


def _add_rect(slide, left, top, width, height, *, fill=LIGHT, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    s.text_frame.text = ""  # placeholder text 제거
    return s


def _add_label_text_in_shape(shape, text, *, size=12, bold=False, color=GRAY,
                              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                              padding=0.1, line_spacing=DEFAULT_LS):
    text = _split_lines(text)
    if isinstance(text, list):
        text = [_split_lines(t) if isinstance(t, str) else t for t in text]
        # flatten one level (each str item may itself have produced a list)
        flat = []
        for t in text:
            if isinstance(t, list):
                flat.extend(t)
            else:
                flat.append(t)
        text = flat
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(padding)
    tf.margin_top = tf.margin_bottom = Inches(padding * 0.6)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    _apply_paragraph_style(p, align=align, line_spacing=line_spacing)
    if isinstance(text, list):
        for i, line in enumerate(text):
            if i > 0:
                p = tf.add_paragraph()
                _apply_paragraph_style(p, align=align, line_spacing=line_spacing)
            run = p.add_run()
            _set_run(run, line, size=size, bold=bold, color=color)
    else:
        run = p.add_run()
        _set_run(run, text, size=size, bold=bold, color=color)


def _add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = Blank


def _add_page_header(slide, kicker, title, page_num=None, total=None):
    # 상단 색 띠
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    # 카테고리 라벨
    _add_text(slide, MARGIN, Inches(0.32), Inches(8), Inches(0.3),
              kicker, size=10, bold=True, color=ACCENT)
    # 타이틀
    _add_text(slide, MARGIN, Inches(0.6), Inches(12), Inches(0.7),
              title, size=26, bold=True, color=NAVY)
    # 가로 구분선
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       MARGIN, Inches(1.35),
                                       SLIDE_W - MARGIN, Inches(1.35))
    line.line.color.rgb = LINE
    line.line.width = Pt(0.75)


# ── 슬라이드 빌더 ────────────────────────────────────────────────────────
def add_title_slide(prs, kicker, title, subtitle):
    s = _add_blank_slide(prs)
    # 좌측 강조 띠
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    _add_text(s, Inches(1.1), Inches(2.4), Inches(11), Inches(0.4),
              kicker, size=14, bold=True, color=ACCENT)
    _add_text(s, Inches(1.1), Inches(2.85), Inches(11.5), Inches(1.4),
              title, size=44, bold=True, color=NAVY)
    _add_text(s, Inches(1.1), Inches(4.3), Inches(11), Inches(0.6),
              subtitle, size=18, color=GRAY)
    _add_text(s, Inches(1.1), Inches(6.6), Inches(11), Inches(0.4),
              "D:\\da\\pilot · v1 (MVP)", size=10, color=GRAY_L)
    return s


def add_section_divider(prs, num, title, body):
    s = _add_blank_slide(prs)
    # 큰 배경 박스
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    _add_text(s, Inches(1.0), Inches(2.6), Inches(2.0), Inches(0.6),
              f"PART {num}", size=14, bold=True, color=ACCENT)
    _add_text(s, Inches(1.0), Inches(3.1), Inches(11), Inches(1.4),
              title, size=40, bold=True, color=WHITE)
    _add_text(s, Inches(1.0), Inches(4.7), Inches(11), Inches(0.8),
              body, size=16, color=RGBColor(0xCC, 0xD2, 0xDB))
    return s


def add_bullet_slide(prs, kicker, title, bullets, page):
    """bullets: list of (text, level) where level 0 = main, 1 = sub."""
    s = _add_blank_slide(prs)
    _add_page_header(s, kicker, title, page)
    top = Inches(1.7)
    for text, level in bullets:
        if level == 0:
            # bullet dot
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     MARGIN, top + Inches(0.13),
                                     Inches(0.12), Inches(0.12))
            dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
            dot.line.fill.background()
            _add_text(s, MARGIN + Inches(0.3), top, SLIDE_W - MARGIN * 2 - Inches(0.3),
                      Inches(0.6), text, size=15, color=NAVY, bold=True)
            top += Inches(0.55)
        else:
            _add_text(s, MARGIN + Inches(0.6), top, SLIDE_W - MARGIN * 2 - Inches(0.6),
                      Inches(0.6), "— " + text, size=13, color=GRAY)
            top += Inches(0.42)
    return s


def add_quote_slide(prs, kicker, title, quote, caption, page):
    s = _add_blank_slide(prs)
    _add_page_header(s, kicker, title, page)
    box = _add_rect(s, MARGIN, Inches(2.2), SLIDE_W - MARGIN * 2, Inches(2.2),
                    fill=LIGHT)
    _add_label_text_in_shape(box, quote, size=22, bold=True, color=NAVY,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                              padding=0.4)
    _add_text(s, MARGIN, Inches(4.7), SLIDE_W - MARGIN * 2, Inches(0.5),
              caption, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    return s


def add_table_slide(prs, kicker, title, headers, rows, page,
                     col_widths=None, top=Inches(1.7), row_h=Inches(0.55),
                     header_fill=NAVY, font_size=12):
    s = _add_blank_slide(prs)
    _add_page_header(s, kicker, title, page)
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    table_left = MARGIN
    table_w = SLIDE_W - MARGIN * 2
    table_h = row_h * n_rows

    shape = s.shapes.add_table(n_rows, n_cols, table_left, top, table_w, table_h)
    table = shape.table

    # 컬럼 폭
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(table_w * (w / total))

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = Inches(0.12)
        tf.margin_top = tf.margin_bottom = Inches(0.06)
        h_lines = h.split("\n") if "\n" in h else [h]
        for li, line in enumerate(h_lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            _apply_paragraph_style(p, align=PP_ALIGN.LEFT, line_spacing=1.15)
            run = p.add_run()
            _set_run(run, line, size=font_size, bold=True, color=WHITE)

    # 본문
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = tf.margin_bottom = Inches(0.06)
            tf.word_wrap = True
            v_lines = val.split("\n") if "\n" in val else [val]
            for li, line in enumerate(v_lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                _apply_paragraph_style(p, align=PP_ALIGN.LEFT, line_spacing=1.2)
                run = p.add_run()
                _set_run(run, line, size=font_size, color=GRAY)

    return s


def add_pitfall_slide(prs, num, name, ai_says, truth, lesson, page):
    """함정 슬라이드 공용 템플릿: 9장 일관."""
    s = _add_blank_slide(prs)
    _add_page_header(s, f"PART 3 · 함정 {num}", name, page)

    # 함정 번호 배지 (좌측 상단 헤더 우측에 큰 숫자)
    badge = _add_rect(s, SLIDE_W - Inches(2.0), Inches(0.5), Inches(1.4), Inches(0.85),
                      fill=ACCENT)
    _add_label_text_in_shape(badge, f"함정 {num}", size=18, bold=True, color=WHITE,
                              align=PP_ALIGN.CENTER, padding=0.05)

    # AI가 잘못 잡은 것 (좌)
    half_w = (SLIDE_W - MARGIN * 2 - Inches(0.4)) // 2
    box_top = Inches(1.65)
    box_h = Inches(4.1)

    box_l = _add_rect(s, MARGIN, box_top, half_w, box_h, fill=LIGHT)
    _add_text(s, MARGIN + Inches(0.25), box_top + Inches(0.2), half_w - Inches(0.5),
              Inches(0.4), "AI가 잘못 잡은 것", size=12, bold=True, color=ACCENT)
    _add_text(s, MARGIN + Inches(0.25), box_top + Inches(0.7), half_w - Inches(0.5),
              box_h - Inches(0.9), ai_says, size=13, color=GRAY)

    # 진실 (우)
    box_r = _add_rect(s, MARGIN + half_w + Inches(0.4), box_top, half_w, box_h,
                      fill=LIGHT)
    _add_text(s, MARGIN + half_w + Inches(0.65), box_top + Inches(0.2),
              half_w - Inches(0.5), Inches(0.4),
              "진실", size=12, bold=True, color=GREEN)
    _add_text(s, MARGIN + half_w + Inches(0.65), box_top + Inches(0.7),
              half_w - Inches(0.5), box_h - Inches(0.9),
              truth, size=13, color=GRAY)

    # 교훈 (하단 풀폭 띠)
    lesson_top = box_top + box_h + Inches(0.25)
    lesson_box = _add_rect(s, MARGIN, lesson_top, SLIDE_W - MARGIN * 2,
                           Inches(0.85), fill=NAVY)
    _add_label_text_in_shape(lesson_box, "→ " + lesson, size=14, bold=True,
                              color=WHITE, align=PP_ALIGN.LEFT, padding=0.25)
    return s


def add_kod_diagram_slide(prs, page):
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 4 · 핵심 모형", "KOD ↔ 원장 — 정책과 적용의 분리", page)

    # 상단: 정책 본체 (KOD)
    kod_w = Inches(4.0); kod_h = Inches(1.1)
    kod_left = (SLIDE_W - kod_w) // 2
    kod = _add_rect(s, kod_left, Inches(1.85), kod_w, kod_h,
                    fill=NAVY, line=NAVY)
    _add_label_text_in_shape(kod, ["KOD_ITN / KOD_ETN",
                                     "정책의 주인 (한도·수수료·우선순위·할인)"],
                              size=14, bold=True, color=WHITE,
                              align=PP_ALIGN.CENTER, padding=0.15)
    _add_text(s, MARGIN, Inches(1.55), Inches(4), Inches(0.3),
              "1차 영향", size=11, bold=True, color=ACCENT)

    # 중간: 화살표 4개 (KOD → 원장 4종)
    ledgers = [
        ("결제\nIAS", "(Issuer Authorization)"),
        ("충전\nCS",  "(Charge Service)"),
        ("환불\nRS",  "(Refund Service)"),
        ("카드\nCMS", "(Card Management)"),
    ]
    led_w = Inches(2.5); led_h = Inches(1.4)
    gap = Inches(0.35)
    total_w = led_w * 4 + gap * 3
    start = (SLIDE_W - total_w) // 2
    led_top = Inches(4.7)

    for i, (name, sub) in enumerate(ledgers):
        left = start + (led_w + gap) * i
        # 화살표 (KOD 하단 중앙 → 원장 상단 중앙)
        arr = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     kod_left + kod_w // 2,
                                     Inches(1.85) + kod_h,
                                     left + led_w // 2, led_top)
        arr.line.color.rgb = GRAY_L
        arr.line.width = Pt(1.25)
        # 화살촉
        arr_xml = arr.line._get_or_add_ln()
        tail = etree.SubElement(arr_xml, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med"); tail.set("len", "med")

        box = _add_rect(s, left, led_top, led_w, led_h, fill=LIGHT, line=LINE)
        _add_label_text_in_shape(box, [name, sub], size=12, bold=True, color=NAVY,
                                  align=PP_ALIGN.CENTER, padding=0.1)

    _add_text(s, MARGIN, led_top + led_h + Inches(0.1), Inches(4), Inches(0.3),
              "2차 영향 — 정책을 받아 *적용한 결과*만 적재",
              size=11, bold=True, color=ACCENT)

    # 하단 핵심 카피
    _add_text(s, MARGIN, Inches(6.8), SLIDE_W - MARGIN * 2, Inches(0.4),
              "정책은 KOD가 들고 있다. 원장은 그 정책을 받아 적용한 결과를 적재할 뿐이다.",
              size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return s


def add_loop_diagram_slide(prs, page):
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 4 · 운용", "정정 루프 — 학습이 누적되는 방식", page)

    steps = [
        ("요건 파일\n입력",                LIGHT, NAVY),
        ("AI 분석\n(함정 9 적용)",         LIGHT, NAVY),
        ("결과 산출\nimpact_X.md",         LIGHT, NAVY),
        ("사람 검토\n+ 정정",              ACCENT, WHITE),
        ("로그 누적\nimpact_log.jsonl",    LIGHT, NAVY),
        ("일반 원칙은\nCLAUDE.md로 승격",  LIGHT, NAVY),
    ]
    n = len(steps)
    box_w = Inches(1.95); box_h = Inches(1.25)
    gap = Inches(0.15)
    total_w = box_w * n + gap * (n - 1)
    start_left = (SLIDE_W - total_w) // 2
    top = Inches(2.55)

    rects = []
    for i, (text, fill, fg) in enumerate(steps):
        left = start_left + (box_w + gap) * i
        r = _add_rect(s, left, top, box_w, box_h, fill=fill,
                      line=LINE if fill == LIGHT else None)
        _add_label_text_in_shape(r, text, size=12, bold=True, color=fg,
                                  align=PP_ALIGN.CENTER, padding=0.1,
                                  line_spacing=1.2)
        rects.append((left, top, box_w, box_h))

    # 직선 화살표 (앞→다음)
    for i in range(n - 1):
        l1, t1, w1, h1 = rects[i]
        l2, t2, _, _ = rects[i + 1]
        arr = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     l1 + w1, t1 + h1 // 2,
                                     l2, t2 + h1 // 2)
        arr.line.color.rgb = GRAY_L; arr.line.width = Pt(1.25)
        ln = arr.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")

    # 회귀 화살표 라벨
    _add_text(s, MARGIN, top + box_h + Inches(0.55), SLIDE_W - MARGIN * 2,
              Inches(0.5),
              "↻ 다음 분석은 누적된 정정과 승격된 원칙 위에서 출발",
              size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _add_text(s, MARGIN, Inches(6.0), SLIDE_W - MARGIN * 2, Inches(0.9),
              ["AI는 한 번에 정답을 못 낸다.",
               "사람의 정정이 누적되는 이 로그(IMPACT-INDEX)가 시스템 도식을 점점 채워가는 *유일한* 경로다."],
              size=13, color=GRAY, align=PP_ALIGN.CENTER)
    return s


def add_funnel_slide(prs, page):
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 3 · 함정 통합", "9가지 함정 → 단 하나의 질문으로 수렴", page)

    # 큰 인용 박스
    quote = _add_rect(s, Inches(1.6), Inches(1.85), SLIDE_W - Inches(3.2),
                      Inches(1.95), fill=LIGHT)
    _add_label_text_in_shape(quote,
        ["이 컴포넌트가 본 요건에서",
         "변경의 주체 인가, 아니면 흐름에 등장만 하는가?"],
        size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, padding=0.3, line_spacing=1.35)

    # 패턴 라벨 4개 (등장 / 매칭 / 호출 / 통과)
    labels = ["등장한다", "매칭된다", "호출된다", "통과한다"]
    box_w = Inches(2.2); box_h = Inches(0.95)
    gap = Inches(0.25)
    total = box_w * 4 + gap * 3
    start = (SLIDE_W - total) // 2
    top = Inches(4.3)
    for i, lab in enumerate(labels):
        left = start + (box_w + gap) * i
        r = _add_rect(s, left, top, box_w, box_h, fill=ACCENT)
        _add_label_text_in_shape(r, lab, size=15, bold=True, color=WHITE,
                                  align=PP_ALIGN.CENTER, padding=0.1)
    _add_text(s, MARGIN, top + box_h + Inches(0.15), SLIDE_W - MARGIN * 2,
              Inches(0.5),
              "위 네 가지를 *변경의 주체* 로 착각하는 패턴 = 함정 ①~⑨",
              size=12, color=GRAY, align=PP_ALIGN.CENTER)

    _add_text(s, MARGIN, Inches(6.4), SLIDE_W - MARGIN * 2, Inches(0.6),
              "분석 시 9개 필터를 순서대로 통과시키면 false positive가 크게 줄어든다.",
              size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return s


def add_toc_slide(prs, page):
    s = _add_blank_slide(prs)
    _add_page_header(s, "목차", "Agenda", page)

    parts = [
        ("PART 1", "도입",            "한 줄 요약 · 풀려는 문제 · 왜 어렵나",                        "3 — 5"),
        ("PART 2", "도구 구성",        "셋이 분담 · 폴더 구조 · 분석 흐름",                            "6 — 8"),
        ("PART 3", "사례와 함정 9가지", "SP_001 47→19 · 함정 ①~⑨ · 한 가지 질문으로 수렴",            "9 — 19"),
        ("PART 4", "핵심 모형과 한계",  "정책 vs 원장 · KOD↔원장 도식 · 예외 셋 · 정정 루프 · 로드맵", "20 — 26"),
    ]

    top = Inches(1.85)
    for kicker, title, body, pages in parts:
        # 좌측 PART 라벨 박스
        kbox = _add_rect(s, MARGIN, top, Inches(1.5), Inches(1.05),
                         fill=NAVY)
        _add_label_text_in_shape(kbox, kicker, size=14, bold=True, color=WHITE,
                                  align=PP_ALIGN.CENTER, padding=0.1)
        # 우측 본문
        _add_text(s, MARGIN + Inches(1.75), top + Inches(0.05),
                  Inches(8.5), Inches(0.45),
                  title, size=20, bold=True, color=NAVY)
        _add_text(s, MARGIN + Inches(1.75), top + Inches(0.55),
                  Inches(8.5), Inches(0.45),
                  body, size=12, color=GRAY)
        # 페이지 범위
        _add_text(s, SLIDE_W - Inches(2.2), top + Inches(0.3),
                  Inches(1.7), Inches(0.4),
                  pages, size=14, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
        top += Inches(1.25)
    return s


def add_two_col_slide(prs, kicker, title, left_title, left_body,
                       right_title, right_body, page):
    s = _add_blank_slide(prs)
    _add_page_header(s, kicker, title, page)
    half_w = (SLIDE_W - MARGIN * 2 - Inches(0.4)) // 2
    top = Inches(1.7)
    h = Inches(5.0)

    # 좌
    box_l = _add_rect(s, MARGIN, top, half_w, h, fill=LIGHT)
    _add_text(s, MARGIN + Inches(0.3), top + Inches(0.25),
              half_w - Inches(0.6), Inches(0.5),
              left_title, size=14, bold=True, color=ACCENT)
    _add_text(s, MARGIN + Inches(0.3), top + Inches(0.85),
              half_w - Inches(0.6), h - Inches(1.1),
              left_body, size=13, color=GRAY)

    # 우
    box_r = _add_rect(s, MARGIN + half_w + Inches(0.4), top, half_w, h,
                      fill=LIGHT)
    _add_text(s, MARGIN + half_w + Inches(0.7), top + Inches(0.25),
              half_w - Inches(0.6), Inches(0.5),
              right_title, size=14, bold=True, color=GREEN)
    _add_text(s, MARGIN + half_w + Inches(0.7), top + Inches(0.85),
              half_w - Inches(0.6), h - Inches(1.1),
              right_body, size=13, color=GRAY)
    return s


# ── 본문 (사용자 원문) ────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 27

    # --- 1. 표지 ----------------------------------------------------------
    add_title_slide(prs,
        "WORKSHOP · v1 (MVP)",
        "코나카드 영향 분석 도구",
        "새 요건 한 줄 → 손대야 할 컴포넌트 후보를 도출하는 사람·AI 협업 도구")

    # --- 2. 목차 ----------------------------------------------------------
    add_toc_slide(prs, 2)

    # --- 3. 한 줄 요약 ---------------------------------------------------
    add_quote_slide(prs,
        "PART 1 · 도입",
        "한 줄 요약",
        "새 요건 한 문장 →\n컴포넌트 208개 중 손대야 할 후보를 추려주는 도구",
        "MVP — 수동 요청 / 컴포넌트 단위. v2 — API 연동 자동 + 테이블·컬럼·API List 단위로 고도화 예정.",
        3)

    # --- 4. 풀려는 문제 — 구체 예시 --------------------------------------
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 1 · 도입", "풀려는 문제 — 구체 예시 하나", 4)
    _add_text(s, MARGIN, Inches(1.7), SLIDE_W - MARGIN*2, Inches(0.7),
              ["코나카드 시스템은 컴포넌트 208개로 쪼개져 있다.",
               "각 컴포넌트에 이름·설명·분류가 붙어 있고 components.xlsx 에 들어 있다."],
              size=13, color=GRAY)
    req = _add_rect(s, MARGIN + Inches(0.5), Inches(2.85),
                    SLIDE_W - MARGIN*2 - Inches(1.0), Inches(1.7),
                    fill=NAVY)
    _add_label_text_in_shape(req,
        ["요건 예시",
         "",
         "“ 본인인증 한 회원에 한해서 선불카드 1일 충전 한도를",
         " 50만원에서 100만원으로 올려달라. ”"],
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, padding=0.3)
    _add_text(s, MARGIN, Inches(4.95), SLIDE_W - MARGIN*2, Inches(0.6),
              "질문 — 이 한 문장 때문에 컴포넌트 208개 중 어떤 것을 손봐야 하는가?",
              size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _add_text(s, MARGIN, Inches(5.85), SLIDE_W - MARGIN*2, Inches(0.5),
              "이게 이 워크스페이스가 풀려는 단 하나의 문제다.",
              size=13, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

    # --- 5. 왜 어렵나 — 3가지 -------------------------------------------
    add_bullet_slide(prs, "PART 1 · 도입", "왜 어렵나 — 3가지", [
        ("208개를 사람이 다 훑을 수 없다", 0),
        ("한 명이 모든 도메인을 알지도 못한다.", 1),
        ("자동으로 알아낼 방법이 없다", 0),
        ("코드·위키·API 명세서 어디에도 컴포넌트 간 연결 고리가 없다.", 1),
        ("각 컴포넌트 담당자조차 전체 도식을 모른다.", 1),
        ("단어만 매칭하면 망한다", 0),
        ("\"한도\"라는 단어가 설명에 있다고 다 영향은 아니다.", 1),
        ("한도값을 들고 있는 컴포넌트 / 받기만 하는 컴포넌트 / 그냥 한 번 거론된 컴포넌트가 다 있다.", 1),
    ], 5)

    # --- 6. 셋이 분담 (표) -----------------------------------------------
    add_table_slide(prs, "PART 2 · 도구 구성", "어떻게 풀었나 — 셋이 분담한다",
                     ["누가", "뭘 하나"],
                     [["파이썬 스크립트",
                       "208개 카탈로그를 읽기 좋게 정리만 한다.\n점수도 안 매기고 추론도 안 한다."],
                      ["AI",
                       "요건 본문과 정리된 카탈로그를 같이 읽고\n어느 컴포넌트를 건드려야 하는지 추론한다."],
                      ["사람",
                       "AI 결과를 보고 틀린 것을 정정한다.\n정정 내용은 로그에 누적해서 다음 분석에 쓴다."]],
                     6, col_widths=[2, 8], row_h=Inches(0.95), font_size=14)
    s = prs.slides[-1]
    _add_text(s, MARGIN, Inches(5.85), SLIDE_W - MARGIN*2, Inches(1.0),
              ["핵심 — AI는 한 번에 정답을 못 낸다.",
               "그래서 사람이 정정 → 정정 사유가 쌓임 → 다음 분석은 그 위에서 출발 하는 루프로 돌린다."],
              size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # --- 7. 폴더 안에 뭐가 있나 (표) -------------------------------------
    add_table_slide(prs, "PART 2 · 도구 구성", "폴더 안에 뭐가 있나 (D:\\da\\pilot)",
                     ["파일", "뭔가", "누가 만드나"],
                     [["components.xlsx",   "컴포넌트 208개 원본 (이름·설명·분류)",      "사람"],
                      ["xlsx_to_md.py",     "xlsx → 사람용 md 변환 스크립트",            "바이브 코딩"],
                      ["components.md",     "사람이 읽는 카탈로그",                        "AI"],
                      ["impact.py",         "xlsx → AI용 압축 카탈로그 변환 스크립트",   "바이브 코딩"],
                      ["impact_context.md", "AI가 빨리 훑는 압축 카탈로그",               "AI"],
                      ["reqs/<요건>.md",     "분석할 새 요건",                              "사람"],
                      ["impact_<요건>.md",   "분석 결과 (어느 컴포넌트가 영향)",           "AI"],
                      ["impact_log.jsonl",  "분석 호출 누적 로그",                          "AI"],
                      ["CLAUDE.md",         "AI 절차서 + 정정으로 배운 함정 9가지",       "사람"]],
                     7, col_widths=[3, 9, 2], row_h=Inches(0.48), font_size=12)
    s = prs.slides[-1]
    _add_text(s, MARGIN, Inches(6.45), SLIDE_W - MARGIN*2, Inches(0.5),
              "스크립트가 둘인 이유 — 사람이 읽는 카탈로그와 AI가 빨리 훑는 카탈로그는 목적이 달라서 따로 만든다.",
              size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # --- 8. 분석 한 번이 어떻게 흘러가나 ---------------------------------
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 2 · 도구 구성",
                     "분석 한 번이 어떻게 흘러가나", 8)
    flow = [
        ("1", "정보 업데이트 체크",
         "impact_context.md 가 components.xlsx 보다 오래됐으면 python impact.py 로 다시 만든다."),
        ("2", "자료 읽기",
         "AI가 요건 본문 + 압축 카탈로그를 같이 읽는다."),
        ("3", "분류 좁히기",
         "이 요건이 건드릴 만한 영역이 어디? 보수적으로 넓게 잡는다 (대개 8~16개 분류)."),
        ("4", "컴포넌트별 판정",
         "분류 안의 모든 컴포넌트를 하나씩 보면서 “이 컴포넌트 자체를 고쳐야 하는가?”를 묻는다. TOP-N 자르기 안 함."),
        ("5", "저장 + 로그",
         "결과는 impact_<요건>.md , 호출 한 줄은 impact_log.jsonl 에 추가."),
    ]
    top = Inches(1.65)
    for num, title, body in flow:
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, top,
                                     Inches(0.7), Inches(0.7))
        circle.fill.solid(); circle.fill.fore_color.rgb = NAVY
        circle.line.fill.background()
        _add_label_text_in_shape(circle, num, size=18, bold=True, color=WHITE,
                                  align=PP_ALIGN.CENTER, padding=0.05)
        _add_text(s, MARGIN + Inches(1.0), top - Inches(0.02),
                  SLIDE_W - MARGIN*2 - Inches(1.0), Inches(0.45),
                  title, size=15, bold=True, color=NAVY)
        _add_text(s, MARGIN + Inches(1.0), top + Inches(0.42),
                  SLIDE_W - MARGIN*2 - Inches(1.0), Inches(0.7),
                  body, size=12, color=GRAY)
        top += Inches(1.1)

    # --- 9. 사례: SP_001 — 47 → 19 ---------------------------------------
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 3 · 사례와 함정", "왜 한 번에 안 됐나 — SP_001 사례", 9)
    _add_text(s, MARGIN, Inches(1.7), SLIDE_W - MARGIN*2, Inches(0.4),
              "요건: \"본인인증 회원에 한해 선불카드 1일 충전 한도 50→100만원\"",
              size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # 큰 숫자 비교 박스 3개
    box_w = Inches(3.6); box_h = Inches(2.4)
    gap = Inches(0.3)
    total = box_w * 3 + gap * 2
    start = (SLIDE_W - total) // 2
    top = Inches(2.5)

    items = [
        ("47", "1차 분석", "AI가 영향이라고 잡은 컴포넌트", LIGHT, NAVY),
        ("→", "6번 정정 후", "사람이 함정을 짚어 정정 누적", LIGHT, GRAY),
        ("19", "최종", "실제 영향 컴포넌트", LIGHT, GREEN),
    ]
    for i, (big, title, sub, fill, fg) in enumerate(items):
        left = start + (box_w + gap) * i
        b = _add_rect(s, left, top, box_w, box_h, fill=fill, line=LINE)
        _add_text(s, left, top + Inches(0.2), box_w, Inches(0.4),
                  title, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _add_text(s, left, top + Inches(0.55), box_w, Inches(1.4),
                  big, size=72, bold=True, color=fg, align=PP_ALIGN.CENTER)
        _add_text(s, left, top + Inches(1.85), box_w, Inches(0.5),
                  sub, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    bot = _add_rect(s, MARGIN, Inches(5.4), SLIDE_W - MARGIN*2, Inches(0.85),
                    fill=ACCENT)
    _add_label_text_in_shape(bot,
        "→ 28개가 잘못 잡혔다. 무엇을 잘못 봤는지 — 다음 9장에서 함정 ①~⑨로 정리.",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT, padding=0.25)

    # --- 10~18. 함정 ①~⑨ -----------------------------------------------
    add_pitfall_slide(prs, "①",
        "카탈로그 어휘 매칭 — \"단어가 같다 = 영향\" 아님",
        ["대표 사례 — CVS",
         "",
         "카탈로그에 \"충전한도 상향 시 사용\"이라고 적힌 컴포넌트(CVS)가 있다.",
         "AI는 \"본 요건이 정확히 충전한도 상향이네!\" 하고 1순위로 잡았다."],
        ["카탈로그에 그렇게 적힌 이유는",
         "\"신분증 인증이 충전한도 상향의 *방법*\"이기 때문.",
         "",
         "본 요건은 한도값 을 바꾸는 것이지 *인증과는 관계 없음*.",
         "→ CVS는 무관."],
        "단어가 겹쳐도, 그 컴포넌트가 본 요건의 변경 행위 에 해당하는지를 한 단계 더 풀어봐야 한다.",
        10)

    add_pitfall_slide(prs, "②",
        "호출 vs 변경 — \"흐름에 등장 = 영향\" 아님",
        ["대표 사례 — CMS · MAP · DCP",
         "",
         "카드 충전이 일어날 때 CMS(카드 상태), MAP(회원 상태), DCP(모바일 카드 상태)를 호출한다.",
         "→ AI는 \"충전 한도가 바뀌면 이것들도 영향\"이라고 잡았다."],
        ["충전을 위해 각 상태를 *확인하려고* 호출이 일어날 뿐.",
         "",
         "한도값이 바뀌어도 CMS·MAP·DCP의 코드/설정은 바뀌지 않는다.",
         "→ 셋 다 영향 아님."],
        "\"흐름에 등장하는가\"가 아니라 \"이 컴포넌트의 코드/설정을 직접 고치는가\"를 물어야 한다.",
        11)

    add_pitfall_slide(prs, "③",
        "도메인 원장 식별 — \"거래 원장 = 결제든 충전이든 IAS\" 아님",
        ["AI는 \"충전도 거래니까 거래 원장 IAS 영향\"이라고 잡았다."],
        ["도메인마다 원장이 다르다 —",
         " · 결제 = IAS",
         " · 충전 = CS",
         " · 환불 = RS",
         " · 카드 = CMS",
         "",
         "본 요건(충전)의 원장은 CS. IAS는 raw data를 받아 적재만 한다.",
         "한도값이 바뀌어도 IAS 코드는 안 바뀐다."],
        "한 도메인의 정책 변경은 *그 도메인 원장*에만 집중된다. 다른 원장은 raw data 적재만 받는다.",
        12)

    add_pitfall_slide(prs, "④",
        "정책 = KOD, 원장 = 적용자 (가장 중요한 함정)",
        ["AI는 \"충전 한도면 충전 원장 CS 가 정책 본체\"라고 잡았다."],
        ["한도/수수료/우선순위/할인 등 *대부분의 정책*은",
         "KOD_ITN / KOD_ETN 가 들고 있다.",
         "",
         "원장(CS · IAS · RS · CMS)은 KOD 정책을 받아",
         "*적용한 결과*만 자기 DB에 적재한다.",
         "",
         "(예외 — AML=A-Safe(외부) · 이상탐지=FDS/KFDS · 쿠폰=KCPS)"],
        "정책 변경 요건의 1차 영향 은 KOD, 2차 영향 은 원장. 사유에 \"KOD 정책 변경 → 원장 적용\" 관계를 명시.",
        13)

    add_pitfall_slide(prs, "⑤",
        "by-pass 라우팅 — \"흐름이 통과한다 = 영향\" 아님",
        ["대표 사례 — PP · VVAN · EGS · BGS",
         "",
         "충전·결제 흐름이 이들을 통과한다.",
         "→ AI는 \"흐름에 들어 있으니 영향\"이라고 잡았다."],
        ["이들은 API Response를 *그대로 통과시키는* 라우팅·검증 계층.",
         "",
         "정책값이 바뀌어도 자체 코드/설정은 바뀌지 않는다."],
        "\"흐름이 지나간다\"가 아니라 \"라우팅 컴포넌트의 코드/설정 변경이 *필요한가*\"를 따로 물어야 한다.",
        14)

    add_pitfall_slide(prs, "⑥",
        "인증 실행 vs 결과 조회 — \"본인인증 단어 = 영향\" 아님",
        ["대표 사례 — NICE-CORE · NCS · CVS · UIS",
         "",
         "본 요건이 \"본인인증 완료 사용자\"를 조건으로 둠",
         "→ AI는 \"그러면 인증 컴포넌트 다 영향\"이라고 잡았다."],
        ["이들은 본인인증을 *실행*하는 컴포넌트.",
         "",
         "본 요건은 인증 *결과를 조회*만 한다.",
         "→ 인증 실행 자체는 변경 없음."],
        "\"본인인증\" 단어가 등장한다고 인증 컴포넌트를 다 잡지 말 것. 인증을 *수행*하게 만드는지, *결과를 사용*만 하는지 구분.",
        15)

    add_pitfall_slide(prs, "⑦",
        "통계/집계는 원장 데이터 의존 — \"통계가 변할 = 통계 컴포넌트 영향\" 아님",
        ["대표 사례 — SAS · CBSS · TBOSB",
         "",
         "AI는 \"한도가 바뀌면 통계도 바뀌니 통계 컴포넌트 다 영향\"이라고 잡았다."],
        ["이들은 원장 데이터를 받아 *적재만* 한다.",
         "",
         "원장이 잘 적재해 주면 통계는 *자동으로* 새 값이 들어간다.",
         "",
         "→ 통계 변화는 *자연스러운 결과*."],
        "통계 컴포넌트 자체에 *새 지표·차원·집계 로직*이 명시적으로 요건에 포함된 경우에만 영향.",
        16)

    add_pitfall_slide(prs, "⑧",
        "인프라 일반론 — \"변경이 일어난다 = 인프라도 다 영향\" 아님",
        ["대표 사례 — KAFKA · SCC · KMS · AGS · KAS",
         "",
         "AI는 \"API가 추가되니 인프라도 다 영향\"이라고 잡았다."],
        ["이들 인프라는 정책 변경의 *직접 영향*을 받지 않는다.",
         "",
         " · 메시지 흐름이 새로 추가 되거나",
         " · 설정이 해당 컴포넌트가 관리 하는 것이거나",
         " · 키 정책 자체가 변경 되는 경우에만 영향."],
        "그냥 \"API 추가됨\" 정도로 인프라를 다 잡지 말 것.",
        17)

    add_pitfall_slide(prs, "⑨",
        "운영 데스크 짝 — \"KOD_ITN만 잡으면 됨\" 아님",
        ["AI는 한도 정책 본체로 KOD_ITN만 잡고 KOD_ETN을 빠뜨렸다.",
         "(이번엔 false negative — 누락 사례)"],
        ["정책 등록/관리 영향이면",
         "KOD_ITN(내부)와 KOD_ETN(외부)이 짝으로 함께",
         "정책 정보를 제공한다.",
         "",
         "한쪽만 잡으면 누락.",
         "",
         "(다른 짝 쌍 — 챗봇 코어/프론트, 알림 본체/라우팅 등도 동일)"],
        "정책 본체가 등장하면 짝도 같이 검토.",
        18)

    # --- 19. 9가지 → 한 가지 질문 ---------------------------------------
    add_funnel_slide(prs, 19)

    # --- 20. 핵심 모형 — 정책과 원장 분리 -------------------------------
    s = _add_blank_slide(prs)
    _add_page_header(s, "PART 4 · 핵심 모형",
                     "가장 중요한 한 가지 — 정책과 원장의 분리", 20)

    big = _add_rect(s, MARGIN, Inches(1.85), SLIDE_W - MARGIN*2, Inches(1.4),
                    fill=NAVY)
    _add_label_text_in_shape(big,
        "정책은 KOD가 들고 있다.\n원장은 그 정책을 받아 적용한 결과를 적재할 뿐이다.",
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, padding=0.3)

    # 좌/우 비교
    half_w = (SLIDE_W - MARGIN*2 - Inches(0.4)) // 2
    top = Inches(3.6)
    h = Inches(2.7)

    box_l = _add_rect(s, MARGIN, top, half_w, h, fill=LIGHT)
    _add_text(s, MARGIN + Inches(0.3), top + Inches(0.25),
              half_w - Inches(0.6), Inches(0.4),
              "KOD = 정책의 주인", size=14, bold=True, color=ACCENT)
    _add_text(s, MARGIN + Inches(0.3), top + Inches(0.8),
              half_w - Inches(0.6), h - Inches(1.0),
              ["한도 / 수수료 / 우선순위 / 할인 같은 정책이",
               "데이터로 KOD에 들어 있다.",
               "",
               "▸ KOD_ITN — 내부용 운영 데스크",
               "▸ KOD_ETN — 외부용 운영 데스크"],
              size=13, color=GRAY)

    box_r = _add_rect(s, MARGIN + half_w + Inches(0.4), top, half_w, h, fill=LIGHT)
    _add_text(s, MARGIN + half_w + Inches(0.7), top + Inches(0.25),
              half_w - Inches(0.6), Inches(0.4),
              "도메인 원장 = 적용자",
              size=14, bold=True, color=GREEN)
    _add_text(s, MARGIN + half_w + Inches(0.7), top + Inches(0.8),
              half_w - Inches(0.6), h - Inches(1.0),
              ["충전 = CS / 결제 = IAS / 환불 = RS / 카드 = CMS",
               "",
               "정책을 KOD에서 받아 거래에 *적용*하고",
               "결과를 자기 DB에 *적재*만 한다.",
               "",
               "한도 50→100 변경 시 1차 = KOD, 2차 = CS · 흐름 컴포넌트"],
              size=13, color=GRAY)

    _add_text(s, MARGIN, Inches(6.65), SLIDE_W - MARGIN*2, Inches(0.4),
              "이걸 모르면 \"원장이 핵심이지\" 하고 IAS·CMS를 잘못 1순위로 잡게 된다.",
              size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # --- 21. KOD ↔ 원장 도식 ---------------------------------------------
    add_kod_diagram_slide(prs, 21)

    # --- 22. 예외 셋 — AML/이상탐지/쿠폰 (표) ----------------------------
    add_table_slide(prs, "PART 4 · 핵심 모형",
                     "단, 예외 셋 — 정책의 본체가 KOD가 아닌 도메인",
                     ["도메인", "정책 주인", "비고"],
                     [["AML\n(자금세탁방지)",
                       "외부 솔루션 A-Safe\n(우리 시스템 밖)",
                       "AMLS는 A-Safe ↔ 코나 사이 중재 어댑터.\n정책 본체 아님."],
                      ["이상탐지 룰",
                       "FDS / KFDS",
                       "FDS = 기존 룰 수행 (룰 추가 불가)\nKFDS = 신규 룰 추가 + 수행"],
                      ["쿠폰 정책",
                       "KCPS\n(KonaCard 쿠폰 시스템)",
                       "외부/내부 쿠폰 발행·취소·정책을 함께 보유.\nKOD 아님."]],
                     22, col_widths=[3, 4, 6], row_h=Inches(0.95), font_size=12)
    s = prs.slides[-1]
    _add_text(s, MARGIN, Inches(5.7), SLIDE_W - MARGIN*2, Inches(0.5),
              "이 예외도 모두 정정 과정에서 사용자가 알려줘서 새겨진 학습이다.",
              size=12, color=GRAY, align=PP_ALIGN.CENTER, bold=True)

    # --- 23. 정정 루프 도식 ---------------------------------------------
    add_loop_diagram_slide(prs, 23)

    # --- 24. MVP → v2 로드맵 --------------------------------------------
    add_two_col_slide(prs, "PART 4 · 한계와 앞으로",
        "지금(MVP)과 다음(v2)",
        "MVP — 지금",
        ["▸ 입력 — 사람이 요건 .md 파일로 던짐",
         "▸ 단위 — 컴포넌트 208개 단위까지",
         "▸ 호출 — 수동 (CLI)",
         "▸ 학습 — impact_log.jsonl + CLAUDE.md 정정 누적",
         "",
         "한계",
         " · 결과는 비결정적 (호출마다 약간 다름)",
         " · 머티리얼이 오래되면 누락 위험",
         " · 외부 시스템 미참조",
         " · 사람 검토 *필수* — 후보군일 뿐"],
        "v2 — 앞으로",
        ["▸ 입력 — 요건 등록 시스템과 API 연동 (자동 트리거)",
         "▸ 단위 — 컴포넌트 + *테이블 · 컬럼 · API List* 까지 확장",
         "▸ 호출 — 자동 (요건 등록과 동시에)",
         "▸ 학습 — 동일하게 정정 루프 (사람 검토는 계속 필수)",
         "",
         "기대",
         " · 영향 후보를 더 *세밀*하게 (어느 테이블·어느 API)",
         " · 영향 분석을 요건 등록 *워크플로우*에 직접 합류",
         " · 정정 누적 → 점점 도식이 채워지는 효과 지속"],
        24)

    # --- 25. 산출물 형식 -------------------------------------------------
    s = _add_blank_slide(prs)
    _add_page_header(s, "BACKUP", "산출물 형식 — 두 가지", 25)
    half_w = (SLIDE_W - MARGIN*2 - Inches(0.4)) // 2
    top = Inches(1.7); h = Inches(5.0)

    bl = _add_rect(s, MARGIN, top, half_w, h, fill=LIGHT)
    _add_text(s, MARGIN + Inches(0.3), top + Inches(0.25),
              half_w - Inches(0.6), Inches(0.4),
              "impact_<basename>.md (분석 결과)",
              size=13, bold=True, color=ACCENT)
    _add_text(s, MARGIN + Inches(0.3), top + Inches(0.85),
              half_w - Inches(0.6), h - Inches(1.1),
              ["헤더 — 요건 파일 경로 · 분석 시각(ISO) ·",
               " 도출 분류 수 · 영향 컴포넌트 수",
               "",
               "## 도출된 분류",
               " 대 > 중 > 소 트리. 각 분류 도출 사유 1줄",
               "",
               "## 영향 컴포넌트",
               " 분류별 그룹. 표:",
               " Name | Full Name | 영향 사유 | components.md 링크",
               "",
               "## 검토 메모 (선택)",
               " 분류는 도출됐으나 영향 없음으로 판정한 컴포넌트의",
               " 누락 사유"],
              size=11, color=GRAY)

    br = _add_rect(s, MARGIN + half_w + Inches(0.4), top, half_w, h, fill=LIGHT)
    _add_text(s, MARGIN + half_w + Inches(0.7), top + Inches(0.25),
              half_w - Inches(0.6), Inches(0.4),
              "impact_log.jsonl (호출 누적 로그)",
              size=13, bold=True, color=ACCENT)
    _add_text(s, MARGIN + half_w + Inches(0.7), top + Inches(0.85),
              half_w - Inches(0.6), h - Inches(1.1),
              ["한 줄 = 분석 한 번. 스키마 —",
               "",
               "{",
               " \"ts\": \"<ISO>\",",
               " \"requirement_file\": \"<absolute path>\",",
               " \"derived_categories\": [\"대>중>소\", ...],",
               " \"impacted\": [{\"name\":\"IAS\",\"reason\":\"...\"}],",
               " \"summary\": \"<1~2문장 요약>\"",
               "}",
               "",
               "이 누적이 IMPACT-INDEX 의 실체.",
               "정정 라인은 correction_of: <이전 ts>로 연결."],
              size=11, color=GRAY)

    # --- 26. 도구 ---------------------------------------------------------
    add_two_col_slide(prs, "BACKUP", "두 스크립트의 역할",
        "xlsx_to_md.py",
        ["역할 — components.xlsx → 사람용 components.md",
         "",
         "▸ slugify() — 한글/영숫자/하이픈 슬러그",
         "▸ 대/중분류 그룹핑 + TOC + 컴포넌트 상세",
         "",
         "호출 시점 — components.xlsx 갱신 시 수동 재실행",
         "(components.md는 자동 생성물 — 직접 편집 금지)"],
        "impact.py",
        ["역할 — components.xlsx → AI용 impact_context.md",
         "",
         "▸ 분류 트리 + 컴포넌트별 한 줄 요약 (100자 cap)",
         "▸ 결정론적, LLM 호출 없음",
         "",
         "호출 시점 — components.xlsx 갱신 시,",
         "또는 분석 1단계 freshness 확인 후 필요 시"],
        26)

    # --- 27. Q&A / 끝 ----------------------------------------------------
    # (제거됨)

    prs.save(OUT)
    return prs


if __name__ == "__main__":
    prs = build()
    print(f"Saved {OUT} ({len(prs.slides)} slides)")
